import os
from typing import List, Dict, Any, Optional, Tuple
import pandas as pd
from pptx import Presentation

class PPTParser:
    """PPT解析器，用于从PPT文件中提取文本和图片信息"""
    
    def __init__(self, ppt_file: str):
        """初始化解析器
        
        Args:
            ppt_file: PPT文件路径
        """
        self.ppt_file = ppt_file
        self.slide_contents = []
    
    def parse(self) -> List[Dict[str, Any]]:
        """解析PPT文件，提取每一页的内容
        
        Returns:
            包含所有幻灯片内容的列表
        """
        print(f"🔍 正在解析PPT文件: {self.ppt_file}")
        
        try:
            presentation = Presentation(self.ppt_file)
            self.slide_contents = []
            
            for slide_index, slide in enumerate(presentation.slides):
                slide_content = {
                    "slide_number": slide_index + 1,
                    "texts": self._extract_texts(slide),
                    "tables": self._extract_tables(slide),
                    "images": self._extract_images(slide),
                    "notes": self._extract_notes(slide),
                }
                self.slide_contents.append(slide_content)
            
            print(f"✅ 成功解析 {len(self.slide_contents)} 页幻灯片")
            return self.slide_contents
            
        except Exception as e:
            print(f"❌ 解析PPT时出错: {str(e)}")
            return []
    
    def _extract_texts(self, slide) -> List[Dict[str, str]]:
        """提取幻灯片中的文本内容
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            文本内容列表
        """
        texts = []
        
        # 遍历所有形状
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_info = {
                    "text": shape.text,
                    "type": "normal_text"
                }
                
                # 检查是否是标题
                if shape.name.startswith("Title"):
                    text_info["type"] = "title"
                
                texts.append(text_info)
        
        return texts
    
    def _extract_tables(self, slide) -> List[Dict[str, Any]]:
        """提取幻灯片中的表格内容
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            表格内容列表
        """
        tables = []
        
        for shape in slide.shapes:
            if shape.has_table:
                # 获取表格行列数
                rows = shape.table.rows
                cols = shape.table.columns
                
                # 创建表格数据
                table_data = []
                for row in range(len(rows)):
                    row_data = []
                    for col in range(len(cols)):
                        cell = shape.table.cell(row, col)
                        row_data.append(cell.text)
                    table_data.append(row_data)
                
                # 创建DataFrame
                df = pd.DataFrame(table_data)
                # 如果第一行是表头，则使用它作为列名
                if len(df) > 0:
                    df.columns = df.iloc[0]
                    df = df[1:]
                
                tables.append({
                    "dataframe": df,
                    "raw_data": table_data
                })
        
        return tables
    
    def _extract_images(self, slide) -> List[Dict[str, Any]]:
        """提取幻灯片中的图片信息
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            图片信息列表
        """
        images = []
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 图片类型
                image_info = {
                    "type": "image",
                    "description": f"幻灯片中的图片 ({shape.width}, {shape.height})"
                }
                images.append(image_info)
        
        return images
    
    def _extract_notes(self, slide) -> Optional[str]:
        """提取幻灯片的备注
        
        Args:
            slide: 幻灯片对象
            
        Returns:
            备注文本
        """
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            return slide.notes_slide.notes_text_frame.text
        return None
    
    def get_slide_text_content(self, slide_index: int = None) -> str:
        """获取指定幻灯片或所有幻灯片的文本内容
        
        Args:
            slide_index: 幻灯片索引（从1开始），如果为None则返回所有幻灯片内容
            
        Returns:
            幻灯片文本内容
        """
        if not self.slide_contents:
            self.parse()
        
        if slide_index is not None:
            # 获取特定幻灯片的内容
            if 1 <= slide_index <= len(self.slide_contents):
                slide = self.slide_contents[slide_index - 1]
                text_content = ""
                
                # 添加文本内容
                for text_item in slide["texts"]:
                    text_content += f"{text_item['text']}\n\n"
                
                # 添加表格内容
                for i, table in enumerate(slide["tables"], 1):
                    text_content += f"表格 {i}:\n"
                    for row in table["raw_data"]:
                        text_content += " | ".join(row) + "\n"
                    text_content += "\n"
                
                # 添加备注
                if slide["notes"]:
                    text_content += f"备注: {slide['notes']}\n"
                
                return text_content.strip()
            else:
                return f"幻灯片索引超出范围: {slide_index}"
        else:
            # 获取所有幻灯片的内容
            all_content = ""
            for slide in self.slide_contents:
                slide_number = slide["slide_number"]
                all_content += f"==== 幻灯片 {slide_number} ====\n\n"
                
                # 添加文本内容
                for text_item in slide["texts"]:
                    text_type = "标题: " if text_item["type"] == "title" else ""
                    all_content += f"{text_type}{text_item['text']}\n\n"
                
                # 添加表格内容
                for i, table in enumerate(slide["tables"], 1):
                    all_content += f"表格 {i}:\n"
                    for row in table["raw_data"]:
                        all_content += " | ".join(row) + "\n"
                    all_content += "\n"
                
                # 添加图片描述
                for i, image in enumerate(slide["images"], 1):
                    all_content += f"图片 {i}: {image['description']}\n"
                
                # 添加备注
                if slide["notes"]:
                    all_content += f"备注: {slide['notes']}\n"
                
                all_content += "\n" + "="*30 + "\n\n"
            
            return all_content.strip()
    
    def save_as_text(self, output_file: str) -> None:
        """将PPT内容保存为文本文件
        
        Args:
            output_file: 输出文件路径
        """
        content = self.get_slide_text_content()
        
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(content)
        
        print(f"✅ PPT内容已保存到: {output_file}")

def main():
    """测试PPT解析功能"""
    import sys
    
    if len(sys.argv) < 2:
        print("使用方法: python ppt_parser.py <ppt文件路径> [输出文本文件路径]")
        return
    
    ppt_file = sys.argv[1]
    
    if not os.path.exists(ppt_file):
        print(f"错误: 找不到PPT文件 '{ppt_file}'")
        return
    
    # 默认输出文件名
    output_file = os.path.splitext(ppt_file)[0] + ".txt"
    if len(sys.argv) >= 3:
        output_file = sys.argv[2]
    
    # 解析PPT并保存为文本
    parser = PPTParser(ppt_file)
    parser.parse()
    parser.save_as_text(output_file)

if __name__ == "__main__":
    main() 